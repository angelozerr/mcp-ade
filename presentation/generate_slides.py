import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
IMAGES_DIR = os.path.join(SCRIPT_DIR, 'images')

# Colors
BG_DARK = RGBColor(0x1E, 0x1E, 0x2E)
BG_CARD = RGBColor(0x2D, 0x2D, 0x3D)
ACCENT = RGBColor(0x00, 0x7A, 0xCC)
ACCENT2 = RGBColor(0x4E, 0xC9, 0xB0)
TEXT_WHITE = RGBColor(0xE0, 0xE0, 0xE0)
TEXT_GRAY = RGBColor(0xA0, 0xA0, 0xA0)
TEXT_YELLOW = RGBColor(0xDC, 0xDC, 0xAA)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)
RED = RGBColor(0xFF, 0x6B, 0x6B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Override theme hyperlink color (default is 0000FF blue)
for layout in prs.slide_layouts:
    master = layout.slide_master
    for rel in master.part.rels.values():
        if 'theme' in rel.reltype:
            theme_xml = etree.fromstring(rel.target_part.blob)
            for hlink in theme_xml.findall('.//' + qn('a:hlink')):
                for child in list(hlink):
                    hlink.remove(child)
                srgb = etree.SubElement(hlink, qn('a:srgbClr'))
                srgb.set('val', 'E0E0E0')
            for fhlink in theme_xml.findall('.//' + qn('a:folHlink')):
                for child in list(fhlink):
                    fhlink.remove(child)
                srgb = etree.SubElement(fhlink, qn('a:srgbClr'))
                srgb.set('val', 'E0E0E0')
            rel.target_part._blob = etree.tostring(theme_xml, xml_declaration=True, encoding='UTF-8', standalone=True)
    break

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=TEXT_WHITE, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = 'Segoe UI'
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = 'Segoe UI'
    return tf

def add_bullet_list_linked(slide, left, top, width, height, items, font_size=16, color=TEXT_WHITE, bullet_color=ACCENT, link_color=ACCENT):
    """Items are tuples: (link_label, url, description)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (label, url, desc) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.space_before = Pt(2)
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = bullet_color
        run_bullet.font.name = 'Segoe UI'
        run_link = p.add_run()
        run_link.text = label
        run_link.font.size = Pt(font_size)
        run_link.font.name = 'Segoe UI'
        run_link.font.underline = True
        run_link.hyperlink.address = url
        # Force color override — solidFill must come BEFORE hlinkClick in OOXML schema
        rPr = run_link._r.get_or_add_rPr()
        for old in rPr.findall(qn('a:solidFill')):
            rPr.remove(old)
        sf = etree.Element(qn('a:solidFill'))
        c = etree.SubElement(sf, qn('a:srgbClr'))
        c.set('val', '%02X%02X%02X' % (link_color[0], link_color[1], link_color[2]))
        hlink = rPr.find(qn('a:hlinkClick'))
        if hlink is not None:
            rPr.insert(list(rPr).index(hlink), sf)
        else:
            rPr.append(sf)
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(font_size)
        run_desc.font.color.rgb = color
        run_desc.font.name = 'Segoe UI'
    return tf


def add_screenshot(slide, image_name, left=0.8, top=2.2, max_width=11.7, max_height=5.0):
    from PIL import Image
    path = os.path.join(IMAGES_DIR, image_name)
    img = Image.open(path)
    w_px, h_px = img.size
    img.close()
    ratio = w_px / h_px
    w = max_width
    h = w / ratio
    if h > max_height:
        h = max_height
        w = h * ratio
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))

def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def slide_title(slide, title, subtitle=None):
    add_accent_line(slide, 0.8, 0.6, 2)
    add_text_box(slide, 0.8, 0.7, 11, 1, title, font_size=36, bold=True, color=TEXT_WHITE)
    if subtitle:
        add_text_box(slide, 0.8, 1.4, 11, 0.6, subtitle, font_size=18, color=TEXT_GRAY)

# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 1, 1.5, 11, 1.2, "MCP ADE (Agent Development Environment)", font_size=54, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, 5, 2.9, 3.3)
add_text_box(slide, 1, 3.2, 11, 0.8, "Bringing Java tooling to AI assistants\nwithout an IDE", font_size=24, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1, 4.8, 11, 0.5, "Eclipse JDT.LS  +  vscode-java-debug  +  MCP", font_size=24, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1, 5.8, 11, 0.5, "Angelo Zerr", font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2 — What already exists
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "What Already Exists", "MCP servers for code intelligence and debugging")

# Top left: Java-specific MCP
add_card(slide, 0.8, 2.2, 6.0, 2.2)
add_text_box(slide, 1.1, 2.3, 5.5, 0.4, "Java-specific MCP servers", font_size=16, bold=True, color=ACCENT2)
add_bullet_list_linked(slide, 1.1, 2.7, 5.5, 1.5, [
    ("javalens-mcp", "https://github.com/pzalutski-pixel/javalens-mcp", " — 56 tools, Eclipse JDT Core direct"),
    ("codebase-memory-mcp", "https://github.com/DeusData/codebase-memory-mcp", " — tree-sitter + JDT semantic"),
    ("LSP4J-MCP", "https://github.com/stephanj/LSP4J-MCP", " — JDT.LS wrapper via LSP4J"),
    ("java-jdtls-mcp-server", "https://github.com/SachieWang/java-jdtls-mcp-server", " — JDT.LS bridge + AI skills"),
], font_size=12, bullet_color=ACCENT2, link_color=TEXT_WHITE)

# Top right: IDE-based MCP
add_card(slide, 7.0, 2.2, 5.5, 2.2)
add_text_box(slide, 7.3, 2.3, 5, 0.4, "IDE-based MCP servers", font_size=16, bold=True, color=RGBColor(0x56, 0x9C, 0xD6))
add_bullet_list_linked(slide, 7.3, 2.7, 5, 1.5, [
    ("JetBrains MCP Server", "https://www.jetbrains.com/help/rider/mcp-server.html", " — built-in since 2025.2 (IntelliJ, PyCharm...)"),
    ("vscode-mcp-server", "https://github.com/juehang/vscode-mcp-server", " — VS Code editing features via MCP"),
], font_size=12, bullet_color=RGBColor(0x56, 0x9C, 0xD6), link_color=TEXT_WHITE)
add_text_box(slide, 7.3, 3.6, 5, 0.4, "Powerful, but require a running IDE", font_size=12, color=TEXT_GRAY)

# Bottom left: MCP + LSP only
add_card(slide, 0.8, 4.6, 6.0, 1.2)
add_text_box(slide, 1.1, 4.65, 5.5, 0.35, "MCP + LSP only", font_size=16, bold=True, color=ACCENT2)
add_bullet_list_linked(slide, 1.1, 5.0, 5.5, 0.7, [
    ("lsp-mcp", "https://github.com/Tritlo/lsp-mcp", " — hover, completion, code actions"),
    ("mcp-language-server", "https://github.com/isaacphi/mcp-language-server", " — definition, references, rename"),
], font_size=12, bullet_color=TEXT_GRAY, link_color=TEXT_WHITE)

# Bottom right: MCP + DAP only
add_card(slide, 7.0, 4.6, 5.5, 1.2)
add_text_box(slide, 7.3, 4.65, 5, 0.35, "MCP + DAP only", font_size=16, bold=True, color=ORANGE)
add_bullet_list_linked(slide, 7.3, 5.0, 5, 0.7, [
    ("dap-mcp", "https://github.com/KashunCheng/dap_mcp", " — breakpoints, step, evaluate"),
    ("mcp-debug-server", "https://github.com/JesseObrien/mcp-debug-server", " — DAP integration for Claude"),
], font_size=12, bullet_color=TEXT_GRAY, link_color=TEXT_WHITE)

# Summary card at the bottom
add_card(slide, 0.8, 6.0, 11.7, 1.4)
add_text_box(slide, 1.2, 6.05, 11, 0.35, "Good tools, but...", font_size=17, bold=True, color=ORANGE)
add_bullet_list(slide, 1.2, 6.4, 11, 0.9, [
    "LSP or DAP — never both together  ·  No server collaboration  ·  No multi-language platform",
    "No Quarkus/Qute/MicroProfile (no JDT delegate command handler)",
    "IDE-based MCP: powerful but require a running IDE — not usable in CI, containers, or headless agents",
], font_size=13, bullet_color=ORANGE)

# ============================================================
# SLIDE 3 — Beyond LSP: Why Advanced Java Tooling is Hard
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Beyond LSP", "LSP features alone are not enough for robust Java MCP tools")

# Left: Two existing approaches
add_card(slide, 0.8, 2.2, 5.8, 2.2)
add_text_box(slide, 1.1, 2.3, 5.4, 0.4, "Approach 1: JDT Core direct (javalens-mcp)", font_size=15, bold=True, color=ACCENT2)
add_bullet_list(slide, 1.1, 2.8, 5.4, 1.4, [
    "56 advanced Java tools (type hierarchy, refactoring...)",
    "Full JDT API access — robust results",
    "But: no Quarkus/Qute/Liberty, no server collaboration",
    "Single server, no delegate command handler support",
], font_size=13, bullet_color=ACCENT2)

add_card(slide, 0.8, 4.7, 5.8, 2.2)
add_text_box(slide, 1.1, 4.8, 5.4, 0.4, "Approach 2: JDT.LS via LSP protocol only (LSP4J-MCP...)", font_size=15, bold=True, color=TEXT_GRAY)
add_bullet_list(slide, 1.1, 5.3, 5.4, 1.4, [
    "Wraps JDT.LS but only uses LSP protocol (completion, hover...)",
    "Ignores JDT.LS plugins and delegate command handlers",
    "No advanced analysis, no refactoring, no framework support",
    "Cannot reuse vscode-quarkus/liberty/debug plugins",
], font_size=13, bullet_color=TEXT_GRAY)

# Right: MCP LT solution
add_card(slide, 7.0, 2.2, 5.5, 4.7, RGBColor(0x1A, 0x3A, 0x5C))
add_text_box(slide, 7.3, 2.3, 5, 0.5, "MCP ADE (Agent Development Environment): Best of both", font_size=17, bold=True, color=TEXT_YELLOW)
add_text_box(slide, 7.3, 2.8, 5, 0.7,
    "Platform that orchestrates LSP & DAP servers and manages\ntheir communication (like VS Code + vscode-java for Java,\nbut for any language)",
    font_size=12, color=TEXT_WHITE)

add_card(slide, 7.3, 3.3, 4.9, 1.6)
add_text_box(slide, 7.6, 3.35, 4.4, 0.35, "JDT.LS = VS Code extension ecosystem", font_size=14, bold=True, color=ACCENT2)
add_text_box(slide, 7.6, 3.7, 4.4, 1.1,
    "Using JDT.LS (not JDT Core) means reusing plugins from:\n"
    "→ LSP: vscode-quarkus, vscode-microprofile, vscode-liberty\n"
    "→ DAP: vscode-java-debug (requires JDT.LS)\n"
    "All these plugins work with MCP LT out-of-the-box",
    font_size=12, color=TEXT_GRAY)

add_card(slide, 7.3, 5.1, 4.9, 1.7)
add_text_box(slide, 7.6, 5.15, 4.4, 0.35, "+ Advanced tools via mcp-jdtls plugin", font_size=14, bold=True, color=ORANGE)
add_text_box(slide, 7.6, 5.55, 4.4, 1.1,
    "OSGi bundle deployed into JDT.LS:\n"
    "→ Full JDT API access for 80 advanced MCP tools\n"
    "   (hierarchy, refactoring, code quality, search...)\n"
    "→ Best of both: LSP ecosystem + JDT API power",
    font_size=12, color=TEXT_GRAY)

# ============================================================
# SLIDE 4 — The Solution: MCP as the Bridge
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "The Solution: MCP as the Bridge")

add_text_box(slide, 0.8, 2.0, 11, 0.5,
    "An MCP server that manages LSP and DAP servers together and exposes their capabilities as MCP tools",
    font_size=20, color=TEXT_GRAY)

# Architecture diagram
add_card(slide, 0.8, 2.8, 2.8, 1.2)
add_text_box(slide, 0.8, 2.9, 2.8, 0.5, "AI Assistant", font_size=18, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 0.8, 3.4, 2.8, 0.4, "Claude Code / Bob", font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 3.8, 3.2, 1.2, 0.5, "→ MCP →", font_size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_card(slide, 5.0, 2.8, 3.2, 1.8, RGBColor(0x1A, 0x3A, 0x5C))
add_text_box(slide, 5.0, 2.9, 3.2, 0.5, "MCP ADE (Agent Development Environment)", font_size=15, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 5.0, 3.4, 3.2, 0.4, "110+ MCP Tools", font_size=13, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 5.0, 3.8, 3.2, 0.4, "Workspace & Lifecycle Mgr", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 8.4, 2.85, 0.8, 0.4, "→ LSP", font_size=12, color=ACCENT2, bold=True)
add_text_box(slide, 8.4, 3.45, 0.8, 0.4, "→ LSP", font_size=12, color=ACCENT2, bold=True)
add_text_box(slide, 8.4, 4.05, 0.8, 0.4, "→ DAP", font_size=12, color=ORANGE, bold=True)

add_card(slide, 9.3, 2.65, 3.2, 0.7)
add_text_box(slide, 9.3, 2.7, 3.2, 0.4, "Eclipse JDT.LS", font_size=14, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 9.3, 3.0, 3.2, 0.3, "80 Java tools + mcp-jdtls", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_card(slide, 9.3, 3.5, 3.2, 0.55)
add_text_box(slide, 9.3, 3.55, 3.2, 0.4, "Quarkus / Qute / Liberty LS", font_size=13, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)

add_card(slide, 9.3, 4.2, 3.2, 0.55)
add_text_box(slide, 9.3, 4.25, 3.2, 0.4, "vscode-java-debug", font_size=14, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Key differentiators — 4 compact cards at bottom
add_card(slide, 0.6, 5.2, 3.0, 2.0, RGBColor(0x1A, 0x3A, 0x5C))
add_text_box(slide, 0.8, 5.3, 2.6, 0.35, "LSP + DAP together", font_size=14, bold=True, color=TEXT_YELLOW)
add_text_box(slide, 0.8, 5.7, 2.6, 1.3, "JDT.LS + Quarkus/Qute/Liberty (LSP)\nvscode-java-debug (DAP)\nDAP ↔ LSP coordination", font_size=12, color=TEXT_GRAY)

add_card(slide, 3.8, 5.2, 3.0, 2.0, RGBColor(0x1A, 0x3A, 0x5C))
add_text_box(slide, 4.0, 5.3, 2.6, 0.35, "Server collaboration", font_size=14, bold=True, color=TEXT_YELLOW)
add_text_box(slide, 4.0, 5.7, 2.6, 1.3, "contributes: bundles, classpath,\nbindRequest/Notification\n80 Java tools via mcp-jdtls", font_size=12, color=TEXT_GRAY)

add_card(slide, 7.0, 5.2, 2.8, 2.0, RGBColor(0x1A, 0x3A, 0x5C))
add_text_box(slide, 7.2, 5.3, 2.4, 0.35, "Open platform", font_size=14, bold=True, color=TEXT_YELLOW)
add_text_box(slide, 7.2, 5.7, 2.4, 1.3, "47 language extensions\n52 LSP + 8 DAP servers\nAdmin console", font_size=12, color=TEXT_GRAY)

add_card(slide, 10.0, 5.2, 2.8, 2.0, RGBColor(0x1A, 0x3A, 0x5C))
add_text_box(slide, 10.2, 5.3, 2.4, 0.35, "Practical", font_size=14, bold=True, color=TEXT_YELLOW)
add_text_box(slide, 10.2, 5.7, 2.4, 1.3, "VSCode settings.json reuse\nAuto-install servers\nIDE server reuse (study)", font_size=12, color=TEXT_GRAY)

# ============================================================
# SLIDE 5 — Technology Stack
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Technology Stack", "Built on Quarkus MCP Server")

add_text_box(slide, 0.8, 2.0, 11, 0.6,
    "MCP tools are declared via Quarkus MCP Server annotations — zero boilerplate",
    font_size=20, color=TEXT_GRAY)

# Left: code example
add_card(slide, 0.8, 2.8, 6.5, 4.0, RGBColor(0x1A, 0x1A, 0x2E))
add_text_box(slide, 1.0, 2.9, 6, 0.4, "Declaring an MCP tool (Java)", font_size=13, color=TEXT_GRAY)
add_text_box(slide, 1.0, 3.3, 6, 3.3,
    '@Tool(name = "java_find_field_writes",\n'
    '      description = "Find all write accesses to a field")\n'
    'public CompletableFuture<String> findFieldWrites(\n'
    '        @ToolArg(description = "...") String cwd,\n'
    '        @ToolArg(description = "...") String fileUri,\n'
    '        @ToolArg(description = "...") int line,\n'
    '        @ToolArg(description = "...") int character,\n'
    '        Cancellation cancellation,\n'
    '        Progress progress) { ... }',
    font_size=12, color=RGBColor(0xCE, 0xD4, 0xDA))

# Right: stack components
add_card(slide, 7.8, 2.8, 4.8, 1.0, RGBColor(0x1A, 0x3A, 0x5C))
add_text_box(slide, 8.0, 2.9, 4.4, 0.4, "Quarkus MCP Server", font_size=18, bold=True, color=ACCENT)
add_text_box(slide, 8.0, 3.3, 4.4, 0.35, "@Tool, @ToolArg, Progress, Cancellation", font_size=12, color=TEXT_GRAY)

add_card(slide, 7.8, 4.0, 4.8, 0.7)
add_text_box(slide, 8.0, 4.05, 4.4, 0.35, "Quarkus (runtime)", font_size=16, bold=True, color=ACCENT2)
add_text_box(slide, 8.0, 4.35, 4.4, 0.3, "CDI, native binary, fast startup", font_size=12, color=TEXT_GRAY)

add_card(slide, 7.8, 4.9, 4.8, 0.7)
add_text_box(slide, 8.0, 4.95, 4.4, 0.35, "LSP4J", font_size=16, bold=True, color=ACCENT2)
add_text_box(slide, 8.0, 5.25, 4.4, 0.3, "LSP and DAP protocol clients", font_size=12, color=TEXT_GRAY)

add_card(slide, 7.8, 5.8, 4.8, 0.7)
add_text_box(slide, 8.0, 5.85, 4.4, 0.35, "Admin UI", font_size=16, bold=True, color=ORANGE)
add_text_box(slide, 8.0, 6.15, 4.4, 0.3, "Quarkus web resources + WebSocket", font_size=12, color=TEXT_GRAY)

# ============================================================
# SLIDE 6 — Core Concepts
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Core Concepts", "Adding a language server or debug adapter is simple")

concepts = [
    ("Add an LSP server", ACCENT2,
     "Just provide a server.json:\n"
     "name, command, documentSelector.\n"
     "MCP LT handles lifecycle, workspace isolation,\n"
     "traces, and contribution wiring automatically."),
    ("Add a DAP server", ORANGE,
     "Same simplicity for debug adapters:\n"
     "name, command, debug templates.\n"
     "Link to an LSP server for coordination\n"
     "(e.g., java-debug ↔ JDT.LS). Done."),
    ("Workspace isolation", RGBColor(0x56, 0x9C, 0xD6),
     "Each project gets its own server instances.\n"
     "Created on-demand when an AI agent calls\n"
     "a tool with a cwd parameter.\n"
     "No cross-project interference."),
    ("Server collaboration", ACCENT,
     "Servers communicate automatically:\n"
     "contributes: bundles, classpath,\n"
     "bindRequest/Notification.\n"
     "Declared in server.json — no code needed."),
]

for i, (title, color, desc) in enumerate(concepts):
    col = i % 2
    row = i // 2
    x = 0.6 + col * 6.3
    y = 2.2 + row * 2.6
    add_card(slide, x, y, 5.9, 2.3)
    add_text_box(slide, x + 0.3, y + 0.15, 5.3, 0.4, title, font_size=22, bold=True, color=color)
    add_text_box(slide, x + 0.3, y + 0.65, 5.3, 1.5, desc, font_size=14, color=TEXT_GRAY)

# ============================================================
# SLIDE 7 — Java: 80 specialized tools
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Java: 80 Specialized MCP Tools", "Powered by Eclipse JDT.LS + custom mcp-jdtls plugin (OSGi bundle)")

java_cats = [
    ("Analysis", "12", ["Type hierarchy", "Call hierarchy", "Control flow", "Change impact"]),
    ("Navigation", "11", ["Go-to-definition", "Hover info", "Javadoc", "Symbol info"]),
    ("Code Search", "11", ["Find references", "Find implementations", "Find tests", "Find unused code"]),
    ("Reference Search", "6", ["Method references", "Field writes", "Casts", "Catch blocks"]),
    ("Refactoring", "18", ["Extract method/variable", "Inline method", "Convert to record", "Move type"]),
    ("Code Generation", "5", ["Getters/setters", "Constructors", "equals/hashCode", "toString"]),
    ("Diagnostics", "4", ["Validate syntax", "Quick fixes", "Diagnose & fix"]),
    ("Code Quality", "5", ["Large classes", "Naming violations", "Circular deps", "Possible bugs"]),
    ("Framework", "3", ["HTTP endpoints", "JPA model", "DI registrations"]),
    ("Project", "4", ["Project structure", "Classpath info", "Dependencies"]),
]

for i, (cat, count, items) in enumerate(java_cats):
    col = i % 5
    row = i // 5
    x = 0.6 + col * 2.5
    y = 2.2 + row * 2.6
    add_card(slide, x, y, 2.3, 2.3)
    add_text_box(slide, x + 0.15, y + 0.1, 2, 0.35, f"{cat} ({count})", font_size=13, bold=True, color=ACCENT2)
    for j, item in enumerate(items):
        add_text_box(slide, x + 0.15, y + 0.5 + j * 0.35, 2, 0.3, f"▸ {item}", font_size=11, color=TEXT_WHITE)

# ============================================================
# SLIDE 8 — Quarkus & Qute Tooling
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Quarkus & Qute Tooling", "Framework-aware Java development via JDT delegate command handler")

add_card(slide, 0.8, 2.2, 5.5, 2.5)
add_text_box(slide, 1.2, 2.3, 5, 0.5, "Quarkus LS", font_size=22, bold=True, color=ACCENT2)
add_bullet_list(slide, 1.2, 2.9, 5, 1.6, [
    "application.properties completion & validation",
    "REST endpoint discovery (java_get_http_endpoints)",
    "CDI injection analysis (java_get_di_registrations)",
    "Delegates Java resolution to JDT.LS via bind mechanism",
], font_size=15)

add_card(slide, 7, 2.2, 5.5, 2.5)
add_text_box(slide, 7.4, 2.3, 5, 0.5, "Qute LS", font_size=22, bold=True, color=ACCENT2)
add_bullet_list(slide, 7.4, 2.9, 5, 1.6, [
    "Qute template completion & validation",
    "Type-safe template expressions",
    "Java type resolution via JDT delegate commands",
    "Template ↔ Java type binding",
], font_size=15)

# How it works — MCP LT as orchestrator
add_card(slide, 0.8, 5.0, 11.7, 2.2, RGBColor(0x1A, 0x3A, 0x5C))
add_text_box(slide, 1.2, 5.1, 10, 0.4, "How it works: MCP ADE (Agent Development Environment) manages inter-LSP communication", font_size=16, bold=True, color=ACCENT)

# MCP LT container
add_card(slide, 1.2, 5.6, 11, 1.5, RGBColor(0x22, 0x33, 0x55))
add_text_box(slide, 1.4, 5.6, 10, 0.3, "MCP ADE (Agent Development Environment) (orchestrator)", font_size=11, color=TEXT_GRAY)

# Quarkus/Qute LS box inside
add_card(slide, 1.5, 6.0, 3.0, 0.9)
add_text_box(slide, 1.5, 6.05, 3.0, 0.35, "Quarkus / Qute LS", font_size=13, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 6.4, 3.0, 0.3, "needs Java type info", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# bind arrow
add_text_box(slide, 4.7, 6.1, 3.0, 0.5, "← bindRequest/Notification →", font_size=12, color=TEXT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

# JDT.LS box inside
add_card(slide, 7.9, 6.0, 3.8, 0.9)
add_text_box(slide, 7.9, 6.05, 3.8, 0.35, "Eclipse JDT.LS", font_size=13, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 7.9, 6.4, 3.8, 0.3, "delegate command handler → resolves Java types", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9 — Java Debugging with DAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Java Debugging via vscode-java-debug", "Full debugging experience through MCP tools")

add_card(slide, 0.8, 2.2, 5.5, 4.8)
add_text_box(slide, 1.2, 2.3, 5, 0.5, "24 DAP Tools", font_size=22, bold=True, color=ORANGE)
add_bullet_list(slide, 1.2, 2.9, 5, 4, [
    "start_debugging — launch or attach to JVM",
    "set_breakpoint / remove_breakpoint",
    "step_over, step_in, step_out",
    "get_stack_trace — full call stack",
    "get_variables / get_local_variables",
    "evaluate_expression — runtime evaluation",
    "get_console_output — stdout/stderr",
    "list_threads — multi-thread inspection",
    "get_debug_statistics — session summary",
], font_size=15)

add_card(slide, 7, 2.2, 5.5, 4.8)
add_text_box(slide, 7.4, 2.3, 5, 0.5, "Example: AI-driven debugging", font_size=20, bold=True, color=ACCENT)
code = '''User: "There\'s a NullPointerException
in UserService.createUser, debug it"

Claude Code:
  1. set_breakpoint("UserService.java", 42)
  2. start_debugging("com.app.Main")
  3. get_local_variables()
     → user.email = null  ← root cause
  4. evaluate_expression("user.toString()")
  5. get_stack_trace()
     → sees the full call chain'''
add_text_box(slide, 7.4, 2.9, 5, 4, code, font_size=13, color=TEXT_YELLOW, font_name='Consolas')

# ============================================================
# SLIDE 10 — Demo: Java Tools in Action
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Demo: Java Tools in Action", "Real examples of Claude Code using MCP ADE (Agent Development Environment)")

# Left: Java tools demo — TODO: add screenshot
add_card(slide, 0.8, 2.2, 5.5, 4.8, RGBColor(0x25, 0x25, 0x35))
add_text_box(slide, 1.2, 2.3, 5, 0.5, "Java / MicroProfile validation", font_size=18, bold=True, color=ACCENT2)
add_text_box(slide, 3, 4.2, 3, 0.5, "[ TODO: screenshot ]", font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Right: Debugger demo — TODO: add screenshot
add_card(slide, 7, 2.2, 5.5, 4.8, RGBColor(0x25, 0x25, 0x35))
add_text_box(slide, 7.4, 2.3, 5, 0.5, "Java Debugging session", font_size=18, bold=True, color=ORANGE)
add_text_box(slide, 9, 4.2, 3, 0.5, "[ TODO: screenshot ]", font_size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11 — Server Collaboration: Bind Mechanism
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Server Collaboration: Bind Mechanism", "Language servers communicate with each other — no custom code needed")

# Row 1: Qute → JDT (bundles + bindRequest/Notification)
add_card(slide, 1, 2.5, 3, 1.1)
add_text_box(slide, 1, 2.6, 3, 0.5, "Qute LS", font_size=16, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.05, 3, 0.35, "template type resolution", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 4.2, 2.7, 2, 0.8, "contributes →\nbundles + bindRequest/Notification", font_size=13, color=TEXT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

add_card(slide, 6.2, 2.5, 3, 1.1)
add_text_box(slide, 6.2, 2.6, 3, 0.5, "Eclipse JDT.LS", font_size=16, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 6.2, 3.05, 3, 0.35, "delegate command handler", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Row 2: Quarkus → JDT (bundles)
add_card(slide, 1, 3.85, 3, 0.95)
add_text_box(slide, 1, 3.9, 3, 0.4, "Quarkus LS", font_size=16, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.3, 3, 0.35, "properties + YAML support", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 4.2, 3.95, 2, 0.8, "contributes →\nbundles", font_size=12, color=TEXT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

add_card(slide, 6.2, 3.85, 3, 0.95)
add_text_box(slide, 6.2, 3.9, 3, 0.4, "Eclipse JDT.LS", font_size=16, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 6.2, 4.3, 3, 0.35, "Quarkus JDT extension", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Row 3: Quarkus → MicroProfile LS (classpath)
add_card(slide, 1, 5.05, 3, 0.95)
add_text_box(slide, 1, 5.1, 3, 0.4, "Quarkus LS", font_size=16, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 3, 0.35, "properties + YAML support", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 4.2, 5.15, 2, 0.8, "contributes →\nclasspath", font_size=12, color=TEXT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

add_card(slide, 6.2, 5.05, 3, 0.95)
add_text_box(slide, 6.2, 5.1, 3, 0.4, "MicroProfile LS", font_size=16, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 6.2, 5.5, 3, 0.35, "Quarkus JAR on classpath", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Row 4: lemminx-liberty → LemMinX (classpath)
add_card(slide, 1, 6.25, 3, 0.95)
add_text_box(slide, 1, 6.3, 3, 0.4, "lemminx-liberty", font_size=16, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.7, 3, 0.35, "server.xml support", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 4.2, 6.35, 2, 0.8, "contributes →\nclasspath", font_size=12, color=TEXT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

add_card(slide, 6.2, 6.25, 3, 0.95)
add_text_box(slide, 6.2, 6.3, 3, 0.4, "LemMinX (XML)", font_size=16, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 6.2, 6.7, 3, 0.35, "Liberty JAR on classpath", font_size=11, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Explanation on the right
add_card(slide, 9.8, 2.5, 3, 4.1)
add_text_box(slide, 10.1, 2.6, 2.5, 0.4, "Contribution types", font_size=16, bold=True, color=ACCENT)
add_bullet_list(slide, 10.1, 3.1, 2.5, 3, [
    "bundles: deploy OSGi plugins into JDT.LS",
    "bindRequest/Notification: route custom LSP requests to delegate handler",
    "classpath: extend another server's classpath with JARs",
    "All declared in server.json",
], font_size=12, bullet_color=ACCENT)

# ============================================================
# SLIDE 12 — Open Platform: 18 Extensions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Open Platform: 47 Built-in Extensions", "Easy to add your own language — just 2 JSON files, no code required")

add_card(slide, 0.8, 2.2, 7.5, 2.2)
add_text_box(slide, 1.2, 2.3, 7, 0.5, "52 LSP servers (built-in)", font_size=20, bold=True, color=ACCENT2)
add_text_box(slide, 1.2, 2.9, 7, 1.2,
    "Java (JDT.LS)  ·  JavaScript/TS  ·  Angular  ·  Python (Pyright)  ·  Go (gopls)  ·  Rust  ·  C/C++  ·  C#\n"
    "HTML/CSS/JSON  ·  Vue  ·  Svelte  ·  Bash  ·  Ruby  ·  Scala  ·  Kotlin  ·  Dart  ·  Swift  ·  Zig  ·  Haskell\n"
    "Elixir  ·  Erlang  ·  OCaml  ·  F#  ·  Clojure  ·  PHP  ·  Lua  ·  Perl  ·  R  ·  Julia  ·  Fortran  ·  Ada  ·  Crystal\n"
    "Pascal  ·  Terraform  ·  LaTeX  ·  Markdown  ·  Nix  ·  TOML  ·  Elm  ·  Ansible  ·  XML  ·  YAML  ·  Dockerfile  ·  ...",
    font_size=12, color=TEXT_GRAY)

add_card(slide, 0.8, 4.7, 7.5, 1.6)
add_text_box(slide, 1.2, 4.8, 7, 0.5, "8 DAP servers (built-in)", font_size=20, bold=True, color=ORANGE)
add_text_box(slide, 1.2, 5.3, 7, 0.8,
    "Java (vscode-java-debug)  ·  JavaScript (vscode-js-debug)  ·  Python (debugpy)\n"
    "Go (delve)  ·  C/C++ (codelldb)  ·  Dart  ·  PHP  ·  Dockerfile",
    font_size=14, color=TEXT_GRAY)

add_card(slide, 8.8, 2.2, 3.8, 4.1)
add_text_box(slide, 9.1, 2.3, 3.3, 0.5, "Add your own!", font_size=20, bold=True, color=ACCENT)
add_bullet_list(slide, 9.1, 2.9, 3.3, 3, [
    "server.json — defines the server",
    "installer.json — auto-download",
    "No code needed",
    "Drop files in extensions/ folder",
    "SPI for advanced cases\n(e.g., mcp-jdtls plugin)",
], font_size=14, bullet_color=ACCENT)
add_text_box(slide, 0.8, 6.6, 11.7, 0.5,
    "These are default extensions — adding a new language is as simple as creating a folder with 2 JSON files",
    font_size=16, bold=True, color=TEXT_YELLOW, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13 — Adding a New Language
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Adding a New Language", "server.json + installer.json (auto-install) in the extensions/ folder — zero code")

add_screenshot(slide, 'extension-json.png')

# ============================================================
# SLIDE 14 — Admin Console
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Admin Console", "Real-time monitoring at localhost:7654/admin")

tabs = [
    ("Workspaces", "Monitor opened projects,\nserver status (Running/Error),\nstart/stop/restart,\ncontribution diagrams"),
    ("Servers", "All 52 LSP servers,\nenable/disable, install,\nlanguage filter,\nconfiguration"),
    ("Debuggers", "8 debug adapters,\ninstallation status,\nconfiguration,\ntrace level control"),
    ("MCP", "Connected MCP clients,\nlive protocol traces,\ntool calls & results,\nreal-time search"),
]

for i, (tab, desc) in enumerate(tabs):
    x = 0.6 + i * 3.15
    add_card(slide, x, 2.4, 2.9, 3.2)
    add_text_box(slide, x + 0.2, 2.6, 2.5, 0.4, tab, font_size=20, bold=True, color=ACCENT)
    add_text_box(slide, x + 0.2, 3.2, 2.5, 2.2, desc, font_size=14, color=TEXT_GRAY)

add_bullet_list(slide, 0.8, 6.0, 11, 1.2, [
    "Real-time WebSocket updates — live LSP/DAP/MCP trace viewing",
    "Server collaboration diagrams (vis.js) — visualize bind relationships",
], font_size=16)

# ============================================================
# SLIDE 15 — Admin: Servers
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Admin: Servers", "52 LSP servers — enable/disable, install, language filter, configuration")

add_screenshot(slide, 'admin-servers.png')

# ============================================================
# SLIDE 16 — Admin: Debuggers
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Admin: Debuggers", "8 debug adapters — installation status, configuration, trace level")

add_screenshot(slide, 'admin-debuggers.png')

# ============================================================
# SLIDE 17 — Admin: MCP Traces
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Admin: MCP Traces", "Connected clients, live protocol traces, tool calls & results")

add_screenshot(slide, 'admin-mcp-traces.png')

# ============================================================
# SLIDE 18 — Admin: Server Contributions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Admin: Server Contributions", "Visualize bind relationships between LSP and DAP servers (vis.js)")

add_screenshot(slide, 'admin-contributions.png')

# ============================================================
# SLIDE 19 — Admin: Workspace / Servers
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Admin: Workspace / Servers", "Monitor opened projects, server status, traces, contribution diagrams")

add_screenshot(slide, 'admin-workspace-servers.png')

# ============================================================
# SLIDE 20 — Admin: Workspace / Debuggers
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Admin: Workspace / Debuggers", "Per-workspace debug adapter instances — sessions, traces")

add_screenshot(slide, 'admin-workspace-debuggers.png')

# ============================================================
# SLIDE 21 — File Watcher: Keeping Language Servers in Sync
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "File Watcher: Keeping Language Servers in Sync",
            "When an AI agent creates files, language servers must know about them")

# Problem card
add_card(slide, 0.8, 2.2, 5.5, 2.0, RGBColor(0x3A, 0x1A, 0x1A))
add_text_box(slide, 1.2, 2.3, 5, 0.4, "The Problem", font_size=18, bold=True, color=RED)
add_bullet_list(slide, 1.2, 2.8, 5, 1.3, [
    "AI agent creates Java file on disk",
    "JDT.LS doesn't know about it → ClassNotFoundException",
    "Root cause: workspace/didChangeWatchedFiles was not implemented",
], font_size=14, bullet_color=RED)

# Solution card
add_card(slide, 7.0, 2.2, 5.5, 2.0, RGBColor(0x1A, 0x3A, 0x1A))
add_text_box(slide, 7.4, 2.3, 5, 0.4, "The Solution: 3 Mechanisms", font_size=18, bold=True, color=ACCENT2)
add_bullet_list(slide, 7.4, 2.8, 5, 1.3, [
    "File Watcher — automatic, local workspaces (Java NIO)",
    "refresh_workspace — explicit, like F5 in Eclipse IDE",
    "notify_file_changes — programmatic, works with remote too",
], font_size=14, bullet_color=ACCENT2)

# Architecture diagram at bottom
add_card(slide, 0.8, 4.5, 11.7, 2.8, RGBColor(0x1A, 0x3A, 0x5C))
add_text_box(slide, 1.2, 4.6, 11, 0.35, "How it works", font_size=16, bold=True, color=ACCENT)

# File system box
add_card(slide, 1.2, 5.1, 2.2, 0.8)
add_text_box(slide, 1.2, 5.15, 2.2, 0.35, "File System", font_size=14, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.2, 5.5, 2.2, 0.3, "*.java, pom.xml...", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow
add_text_box(slide, 3.5, 5.2, 1.5, 0.4, "→ events →", font_size=12, color=TEXT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

# File Watcher box
add_card(slide, 5.0, 5.1, 2.5, 0.8)
add_text_box(slide, 5.0, 5.15, 2.5, 0.35, "WorkspaceFileWatcher", font_size=12, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 5.0, 5.5, 2.5, 0.3, "batch 500ms + glob filter", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Arrow
add_text_box(slide, 7.6, 5.2, 1.5, 0.4, "→ LSP →", font_size=12, color=TEXT_YELLOW, bold=True, alignment=PP_ALIGN.CENTER)

# Language Servers box
add_card(slide, 9.0, 5.1, 3.2, 0.8)
add_text_box(slide, 9.0, 5.15, 3.2, 0.35, "Language Servers", font_size=14, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 9.0, 5.5, 3.2, 0.3, "didChangeWatchedFiles", font_size=10, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Key features
add_text_box(slide, 1.2, 6.2, 5, 0.9,
    "Static patterns (server.json) + dynamic (registerCapability)\n"
    "Event queuing when server stopped → replay on restart\n"
    "Configurable excludes (fileWatchers.excludePatterns)",
    font_size=11, color=TEXT_GRAY)

add_text_box(slide, 7.0, 6.2, 5, 0.9,
    "JDT.LS: also needs project.refreshLocal() via mcp-jdtls\n"
    "→ refresh_workspace calls both LSP + server-specific refresh\n"
    "No auto-refresh on restart — explicit action by agent or admin",
    font_size=11, color=TEXT_GRAY)

# ============================================================
# SLIDE 22 — Roadmap  (was 21)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Roadmap", "Scaling to large Java projects and beyond")

add_card(slide, 0.8, 2.2, 5.5, 4.8)
add_text_box(slide, 1.2, 2.4, 5, 0.5, "Known challenge:\nlarge project import", font_size=17, bold=True, color=ORANGE)
add_bullet_list(slide, 1.2, 3.0, 5, 3.8, [
    "Tested on Quarkus repo (huge codebase):",
    "1. File scanning: each importer scans all files\n  → very slow on large repos",
    "2. Gradle conflict: build.gradle takes over pom.xml\n  → Maven modules imported as Gradle",
    "3. M2E import: ran 2h+, never completed\n  → JDT.LS never becomes ready",
], font_size=14)

add_card(slide, 7, 2.2, 5.5, 4.8)
add_text_box(slide, 7.4, 2.4, 5, 0.5, "Roadmap", font_size=20, bold=True, color=ACCENT)
add_bullet_list(slide, 7.4, 3.0, 5, 3.8, [
    "Fix JDT.LS import (scan, Gradle conflict)",
    "Maven artifact caching strategy started\n  → Quarkus repo: 2h+ → 48 min",
    "Partial / incremental M2E import",
    "Connect to running JDT.LS in IDE\n  → reuse already-imported workspace",
    "More language extensions",
], font_size=14, bullet_color=ACCENT)

# ============================================================
# SLIDE 23 — Getting Started  (was 22)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide_title(slide, "Getting Started")

add_card(slide, 0.8, 2.2, 11.5, 2.0)
add_text_box(slide, 1.2, 2.3, 11, 0.4, "1. Clone, build, and start MCP ADE (Agent Development Environment)", font_size=18, bold=True, color=ACCENT)
add_text_box(slide, 1.2, 2.8, 11, 1.0, "git clone https://github.com/angelozerr/mcp-ade\ncd mcp-ade\n./mvnw clean install -DskipTests\n./mvnw quarkus:dev -f dev/pom.xml", font_size=14, color=TEXT_YELLOW, font_name='Consolas')

add_card(slide, 0.8, 4.0, 11.5, 1.8)
add_text_box(slide, 1.2, 4.1, 11, 0.4, "2. Configure your MCP client (Claude Code, Bob...)", font_size=18, bold=True, color=ACCENT)
config_code = '''{
  "mcpServers": {
    "mcp-ade": { "url": "http://localhost:7654/mcp" }
  }
}'''
add_text_box(slide, 1.2, 4.6, 11, 1.2, config_code, font_size=16, color=TEXT_YELLOW, font_name='Consolas')

add_card(slide, 0.8, 6.1, 11.5, 0.8)
add_text_box(slide, 1.2, 6.2, 11, 0.4, "3. Ask your AI assistant to work on Java code!", font_size=18, bold=True, color=ACCENT)


# ============================================================
# SLIDE 24 — Thank You  (was 23)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 1, 2, 11, 1, "Thank You!", font_size=54, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
add_accent_line(slide, 5, 3.2, 3.3)

add_text_box(slide, 1, 3.8, 11, 0.8, "MCP ADE (Agent Development Environment)", font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_bullet_list_linked(slide, 3, 5.0, 7, 0.5, [
    ("github.com/angelozerr/mcp-ade", "https://github.com/angelozerr/mcp-ade", ""),
], font_size=18, bullet_color=BG_DARK, link_color=TEXT_WHITE)
add_text_box(slide, 1, 5.5, 11, 0.5, "Admin Console: localhost:7654/admin", font_size=18, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1, 6.5, 11, 0.5, "Questions?", font_size=24, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(SCRIPT_DIR, "MCP_ADE.pptx")
prs.save(output_path)
print(f"Presentation saved to {output_path}")
